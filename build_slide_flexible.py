from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import itertools
import math

from pptx.oxml.xmlchemy import OxmlElement


# =========================================================
# BULLETS: edit this only to change the process
# =========================================================
BULLET_TEXT = """
Start: START
Information: Information
Action: Action/s
Decision: Decision?
  Path "Decision path 1":
    Action:
      Heading
      Action/s
  Path "Decision path 2":
    Action:
      Heading
      • Bullet Actions
  Path "Decision path 3":
    Action:
      Action/s  
End: END
""".strip("\n")


# =========================================================
# COLORS / STYLES
# =========================================================
# Brand/Theme colors
DARK_TEAL = (12, 71, 110)  # kept for outlines unless specified otherwise
DOC_HEADING_RGB = (0x00, 0x50, 0x77)   # #005077
SUB_HEADING_RGB = (0xED, 0x11, 0x65)   # #ED1165 (reserved)
KEY_BOX_RGB     = (0x92, 0xD0, 0x50)   # #92D050
KEY_TEXT_RGB    = (0x00, 0x00, 0x00)   # #000000
FLOW_TEXT_RGB   = (0x00, 0x00, 0x00)   # #000000
ARROW_RGB       = (0x44, 0x72, 0xC4)   # #4472C4
FOOTER_RGB      = (0xBF, 0xBF, 0xBF)   # #BFBFBF

BOX_STYLES = {
    "Start": {
        "fill_rgb": (0x92, 0xD0, 0x50),   # Start lozenge #92D050
        "shape": "rounded",
        "bold_first_line": True,
        "align": "center",
    },
    "Information": {
        "fill_rgb": (0x2F, 0xC9, 0xFF),   # #2FC9FF
        "shape": "rounded",
        "bold_first_line": False,
        "align": "center",
    },
    "Action": {
        "fill_rgb": (0x9D, 0xC3, 0xE6),   # #9DC3E6
        "shape": "rounded",
        "bold_first_line": False,
        "align": "center",
    },
    "Decision": {
        "fill_rgb": (0xEA, 0xB0, 0xFA),   # #EAB0FA
        "shape": "diamond",
        "bold_first_line": True,
        "align": "center",
    },
    "End": {
        "fill_rgb": (0xFF, 0xC0, 0x00),   # End lozenge #FFC000
        "shape": "rounded",
        "bold_first_line": True,
        "align": "center",
    },
}

# =========================================================
# STEP 1. Parse flat bullet list -> HUMAN_FLOW
# =========================================================
def parse_bullets_to_human_flow(bullets_str):
    """
    Expects bullets in this format:

    Start: START
    Information: Information
    Action: Action/s
    Decision: Decision?
      Path "Decision path 1":
        Action:
          Heading
          Action/s
      Path "Decision path 2":
        Action:
          Heading
          • Bullet Actions
      Path "Decision path 3":
        Action:
          Action/s
    End: END
    """

    # Tokenize: capture indent and text
    lines = []
    for raw_line in bullets_str.splitlines():
        if not raw_line.strip():
            continue
        indent = len(raw_line) - len(raw_line.lstrip(" "))
        content = raw_line.strip()
        lines.append((indent, content))

    spine_nodes = []
    i = 0
    while i < len(lines):
        indent, content = lines[i]

        # top-level step like "Start: START" or "Decision: Decision?"
        if indent == 0 and not content.lower().startswith("path "):
            if ":" not in content:
                raise ValueError(f"Cannot parse line: {content}")
            node_type, node_text = content.split(":", 1)
            node_type = node_type.strip()
            node_text = node_text.strip()

            node_entry = {
                "type": node_type,
                "text": [node_text] if node_text else [node_type],
            }

            # Decision is special: it has branches underneath
            if node_type.lower() == "decision":
                decision_paths = []
                i += 1
                while i < len(lines):
                    indent2, content2 = lines[i]

                    # if we hit another indent==0, that means the Decision block ended
                    if indent2 == 0:
                        break

                    # Path "Decision path X":
                    if content2.lower().startswith("path "):
                        label_part = content2[len("Path "):].strip()
                        if not label_part.endswith(":"):
                            raise ValueError(f"Missing ':' after path label in '{content2}'")
                        label_part = label_part[:-1].strip()  # strip trailing ':'

                        # strip quotes if present
                        if ((label_part.startswith('"') and label_part.endswith('"')) or
                            (label_part.startswith("'") and label_part.endswith("'"))):
                            label_text = label_part[1:-1]
                        else:
                            label_text = label_part

                        # gather branch steps under this Path
                        branch_steps = []
                        i += 1
                        while i < len(lines):
                            indent3, content3 = lines[i]
                            if indent3 <= indent2:
                                break

                            # branch steps can be:
                            #   Action:
                            #     line1
                            #     line2
                            # or  End: END
                            if ":" not in content3:
                                raise ValueError(f"Cannot parse branch step line: {content3}")
                            step_type, after_colon = content3.split(":", 1)
                            step_type = step_type.strip()
                            after_colon = after_colon.strip()

                            if after_colon:
                                # e.g. End: END  (single-line)
                                branch_steps.append({
                                    "type": step_type,
                                    "text": [after_colon],
                                })
                                i += 1
                                continue

                            # e.g. Action:  (multiline block follows)
                            # consume deeper-indented lines as text lines of the same box
                            text_lines = []
                            i += 1
                            while i < len(lines):
                                indent4, content4 = lines[i]
                                if indent4 <= indent3:
                                    break
                                text_lines.append(content4)
                                i += 1

                            branch_steps.append({
                                "type": step_type,
                                "text": text_lines if text_lines else [step_type],
                            })

                        decision_paths.append({
                            "label": label_text,
                            "steps": branch_steps,
                        })

                        continue
                    else:
                        raise ValueError(
                            f"Unexpected line under Decision: '{content2}'. "
                            f"Expected Path \"...\":"
                        )

                node_entry["paths"] = decision_paths
                spine_nodes.append(node_entry)
                continue  # don't i += 1 here, we've already advanced

            else:
                # normal spine node like Start/Information/Action/End
                spine_nodes.append(node_entry)
                i += 1
                continue

        else:
            # If we see indentation at top level that isn't a Decision branch, it's invalid
            raise ValueError(
                f"Unexpected indentation at line '{content}'. "
                "Only lines under a Decision's Paths should be indented."
            )

    # Turn spine_nodes list into nested HUMAN_FLOW via .children[0]
    def chainify(nodes_list):
        if not nodes_list:
            return None
        head = nodes_list[0]
        out = {
            "type": head["type"],
            "text": head["text"],
        }
        if "paths" in head:
            out["paths"] = head["paths"]
        tail = chainify(nodes_list[1:])
        if tail:
            out["children"] = [tail]
        return out

    return chainify(spine_nodes)


# =========================================================
# STEP 2. HUMAN_FLOW -> FLOW_SPEC
# =========================================================
def linear_chain_to_spec(steps):
    """
    steps = [
      {"type":"Action","text":["Heading","Action/s"]},
      {"type":"End","text":["END"]}
    ]
    -> nested 'paths': {'next': ...}
    """
    if not steps:
        return None

    head = steps[0]
    node_spec = {
        "type": head["type"],
        "text": head["text"],
    }

    tail_spec = linear_chain_to_spec(steps[1:])
    if tail_spec:
        node_spec["paths"] = {"next": tail_spec}

    return node_spec


def decision_paths_to_spec(decision_node):
    """
    decision_node['paths'] = [
      {'label': 'Decision path 1', 'steps': [...]},
      {'label': 'Decision path 2', 'steps': [...]},
      {'label': 'Decision path 3', 'steps': [...]},
    ]
    -> {
         'path1': {'label':'Decision path 1', 'node': <chain>},
         'path2': {...},
         'path3': {...}
       }
    """
    out = {}
    for idx, branch in enumerate(decision_node["paths"], start=1):
        out[f"path{idx}"] = {
            "label": branch["label"],
            "node": linear_chain_to_spec(branch["steps"])
        }
    return out


def to_flow_spec(human_flow):
    spec = {
        "type": human_flow["type"],
        "text": human_flow["text"],
    }

    if human_flow["type"].lower() == "decision":
        spec["paths"] = decision_paths_to_spec(human_flow)
        return spec

    kids = human_flow.get("children", [])
    if kids:
        spec["paths"] = {"next": to_flow_spec(kids[0])}

    return spec


# =========================================================
# STEP 3. FLOW_SPEC -> positioned NODES + CONNECTORS
# (single END node; branches all route into it)
# =========================================================
# COMPLETE REPLACEMENT for the section from "# 3. Force END node..." through "return nodes, connectors"
# This should replace lines approximately 295-415 in your original file

# COMPLETE REPLACEMENT for add_block_arrow in your build_slide.py
# Replace the existing add_block_arrow function (around line 75-100) with these two functions:
'''
def add_block_arrow(slide, x1, y1, x2, y2, thickness_in=0.08):
    """
    Draw a RIGHT_ARROW from (x1, y1) to (x2, y2).
    Arrow points naturally in the direction from start to end.
    
    Args:
        x1, y1: Start point in inches
        x2, y2: End point in inches  
        thickness_in: Arrow thickness in inches
    """
    
    dx = x2 - x1
    dy = y2 - y1
    
    # If points are the same, don't draw
    if dx == 0 and dy == 0:
        return None
    
    # Calculate actual length and angle (no snapping)
    length = math.sqrt(dx*dx + dy*dy)
    angle_deg = math.degrees(math.atan2(dy, dx))
    
    # Midpoint between start and end
    mid_x = (x1 + x2) / 2.0
    mid_y = (y1 + y2) / 2.0
    
    # Create arrow centered at midpoint
    # Arrow shape is horizontal by default, pointing right
    left = Inches(mid_x - length/2.0)
    top = Inches(mid_y - thickness_in/2.0)
    width = Inches(length)
    height = Inches(thickness_in)
    
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        left, top, width, height
    )
    
    # Style
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0, 112, 192)
    shp.line.width = Pt(0)
    
    # Rotate to actual angle (no snapping)
    shp.rotation = angle_deg
    
    return shp

'''
def add_simple_arrow(slide, x1, y1, x2, y2):
    """
    Draw a simple straight line with arrowhead from (x1, y1) to (x2, y2).
    Matches the clean arrow style in your reference image.
    
    Args:
        x1, y1: Start point in inches
        x2, y2: End point in inches
    """
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches
    
    # Create a straight connector
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    
    # Style the line
    connector.line.color.rgb = RGBColor(0, 0, 139)  # Dark blue
    connector.line.width = Pt(1.5)
    
    # Add arrowhead at the end
    connector.line.end_arrow_type = 2  # Arrow
    
    return connector


def add_block_arrow(slide, x1, y1, x2, y2, thickness_in=0.08):
    """
    Draw a simple straight arrow.
    """
    return add_simple_arrow(slide, x1, y1, x2, y2)

# Enhanced version with flexible layout capabilities

def calculate_box_size(text_lines, node_type, font_size_pt=9):
    """Calculate optimal box width and height based on text content."""
    char_width = 0.06
    max_line_length = max(len(line) for line in text_lines) if text_lines else 10
    base_width = max_line_length * char_width
    
    if node_type == "Decision":
        padding = 0.6
        min_width = 1.2
        max_width = 3.5
    elif node_type == "Start" or node_type == "End":
        padding = 0.4
        min_width = 0.8
        max_width = 2.0
    else:
        padding = 0.5
        min_width = 1.0
        max_width = 3.0
    
    width = min(max(base_width + padding, min_width), max_width)
    
    line_height = 0.18
    num_lines = len(text_lines)
    base_height = num_lines * line_height
    
    if node_type == "Decision":
        vertical_padding = 0.35
        min_height = 0.6
    elif node_type == "Start" or node_type == "End":
        vertical_padding = 0.25
        min_height = 0.3
    else:
        vertical_padding = 0.3
        min_height = 0.35
    
    height = max(base_height + vertical_padding, min_height)
    
    if node_type == "Decision":
        if width < height * 1.3:
            width = height * 1.3
    
    return round(width, 2), round(height, 2)


def add_arrow_connector(slide,
                        x1, y1, x2, y2,
                        line_width_pt=1,
                        rgb=(0, 0, 0),
                        arrow_at='end',
                        label_text=None,
                        prefer='v-first'):
    """
    Draw an ORTHOGONAL (horizontal/vertical only) connector from (x1,y1) to (x2,y2).
    Implemented as 1 or 2 straight segments so there are no diagonal lines.

    slide: pptx.slide.Slide object
    x1, y1, x2, y2: all in EMU (pass Inches(...) at callsite)
    line_width_pt: float, stroke width in points (1 = thin)
    rgb: (R,G,B)
    arrow_at: 'start' | 'end' | 'both'
    label_text: Optional text centered on the longer segment
    prefer: 'h-first' or 'v-first' for L-shaped routing
    """

    def _style_line(conn, with_head=False, with_tail=False):
        conn.line.fill.solid()
        conn.line.fill.fore_color.rgb = RGBColor(*rgb)
        conn.line.width = Pt(line_width_pt)
        ln = conn.line._get_or_add_ln()
        # rounded line caps and joins for smoother look
        try:
            ln.set('cap', 'rnd')  # rounded end caps
            round_join = OxmlElement('a:round')
            ln.append(round_join)
        except Exception:
            pass
        if with_head:
            elt = OxmlElement('a:headEnd')
            elt.set('type', 'arrow')
            elt.set('w', 'sm')
            elt.set('len', 'sm')
            ln.append(elt)
        if with_tail:
            elt = OxmlElement('a:tailEnd')
            elt.set('type', 'arrow')
            elt.set('w', 'sm')
            elt.set('len', 'sm')
            ln.append(elt)

    # If already horizontal or vertical, draw a single straight connector
    if x1 == x2 or y1 == y2:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        _style_line(conn,
                    with_head=(arrow_at in ('start', 'both')),
                    with_tail=(arrow_at in ('end', 'both')))
        target_for_label = conn
        mid_x, mid_y = ( (x1+x2)/2, (y1+y2)/2 )
    else:
        # Two-segment orthogonal route. Default: vertical-first then horizontal.
        if prefer == 'v-first':
            # Segment 1: vertical to dest Y
            seg1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x1, y2)
            _style_line(seg1,
                        with_head=(arrow_at in ('start','both')),
                        with_tail=False)
            # Segment 2: horizontal to dest X (arrowhead here)
            seg2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y2, x2, y2)
            _style_line(seg2,
                        with_head=False,
                        with_tail=(arrow_at in ('end','both')))
            target_for_label = seg2
            mid_x, mid_y = ((x1 + x2)/2, y2)
        else:
            # Horizontal-first then vertical
            seg1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y1)
            _style_line(seg1,
                        with_head=(arrow_at in ('start','both')),
                        with_tail=False)
            seg2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x2, y1, x2, y2)
            _style_line(seg2,
                        with_head=False,
                        with_tail=(arrow_at in ('end','both')))
            target_for_label = seg1
            mid_x, mid_y = ((x1 + x2)/2, y1)

    # Label textbox (transparent)
    if label_text:
        text_width = Inches(1.2)
        text_height = Inches(0.25)
        text_box = slide.shapes.add_textbox(
            mid_x - text_width/2,
            mid_y - text_height/2,
            text_width,
            text_height
        )
        try:
            text_box.shadow.inherit = False
            text_box.shadow.visible = False
        except Exception:
            pass
        text_frame = text_box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label_text
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(*FLOW_TEXT_RGB)
        text_box.fill.background()
        text_box.line.fill.background()

    # Not returning a particular connector because we might create two
    return None


def generate_nodes_connectors_from_flow_flexible(flow_spec):
    """
    Generate nodes and connectors with better fitting layout.
    Uses the ORIGINAL working layout but more carefully tuned.
    """
   

    # Layout constants (inches) - ADJUSTED for better arrow visibility
    spine_x = 6.5  # Center of 13.333" slide
    first_y = 1.5  # Start a bit higher
    y_step = 0.65  # MORE vertical spacing (was 0.5)
    
    branch_y = 3.8  # Decision branch level (adjusted)
    left_branch_x = 2.5   # Left branch position
    right_branch_x = 10.0  # Right branch position (moved further right)
    bottom_action_y = 5.2  # Bottom branch
    
    end_x = 6.5  # Center bottom
    end_y = 6.0  # Moved down slightly (was 5.6)
    
    nodes = []
    connectors = []
    counter = itertools.count(1)
    
    memo = {}
    node_geom = {}
    
    def new_id(prefix):
        return f"{prefix}_{next(counter)}"
    
    def build_spine(spec, x, start_y):
        out = []
        cur = spec
        cur_y = start_y
        while cur:
            lk = id(cur)
            if lk in memo:
                node_id = memo[lk]
                w = node_geom[lk]["w"]
                h = node_geom[lk]["h"]
            else:
                node_id = new_id(cur["type"].lower())
                memo[lk] = node_id
                w, h = calculate_box_size(cur["text"], cur["type"])
                
                # CENTER the box horizontally at x
                nodes.append({
                    "id": node_id,
                    "type": cur["type"],
                    "lines": cur["text"],
                    "left": x - w/2,  # CENTER
                    "top": cur_y,
                    "width": w,
                    "height": h,
                })
                node_geom[lk] = {"w": w, "h": h}
            
            out.append((node_id, cur, x, cur_y, w, h))
            
            nxt = cur.get("paths", {}).get("next")
            if not nxt:
                break
            
            cur_y += y_step
            if nxt["type"] == "Decision":
                cur_y += 0.15  # Extra space before decision
            cur = nxt
        
        return out
    
    # 1. Build the main spine
    spine_chain = build_spine(flow_spec, spine_x, first_y)
    
    # 2. Identify Decision + End nodes
    decision_tuple = None
    end_tuple = None
    for tup in spine_chain:
        node_id, spec_obj, x, y, w, h = tup
        if spec_obj["type"] == "Decision":
            decision_tuple = tup
        if spec_obj["type"] == "End":
            end_tuple = tup
    
    # 3. Force END node to bottom-center coordinates
    final_end_node_id = None
    
    if end_tuple:
        end_node_id, end_spec_obj, _, _, _, _ = end_tuple
        final_end_node_id = end_node_id
        end_w, end_h = calculate_box_size(end_spec_obj["text"], "End")
        for nd in nodes:
            if nd["id"] == end_node_id:
                nd["left"] = end_x - end_w/2  # CENTER
                nd["top"] = end_y
                nd["width"] = end_w
                nd["height"] = end_h
                break
    else:
        # fallback if End missing
        end_spec_obj = {"type": "End", "text": ["END"]}
        lk = id(end_spec_obj)
        end_node_id = new_id("end")
        final_end_node_id = end_node_id
        memo[lk] = end_node_id
        end_w, end_h = calculate_box_size(end_spec_obj["text"], "End")
        nodes.append({
            "id": end_node_id,
            "type": "End",
            "lines": end_spec_obj["text"],
            "left": end_x - end_w/2,  # CENTER
            "top": end_y,
            "width": end_w,
            "height": end_h,
        })
        node_geom[lk] = {"w": end_w, "h": end_h}
    
    # 4. Add arrows between spine nodes
    for (nid_a, spec_a, ax, ay, aw, ah), (nid_b, spec_b, bx, by, bw, bh) in zip(spine_chain, spine_chain[1:]):
        if spec_a["type"] == "Decision" and spec_b["type"] == "End":
            continue
        connectors.append({
            "from_id": nid_a,
            "to_id": nid_b,
            "style": "straight",
            "from_anchor": "bottom",
            "to_anchor": "top",
        })
    
    # 5. If there's no decision node, we're done
    if decision_tuple is None:
        return nodes, connectors
    
    decision_id, decision_spec, decision_x, decision_y, decision_w, decision_h = decision_tuple
    decision_paths = decision_spec.get("paths", {})
    
    def place_branch(branch_key):
        branch_def = decision_paths[branch_key]
        branch_root = branch_def["node"]
        label_text = branch_def["label"]
        
        if branch_key == "path1":  # left box
            bx, by = (left_branch_x, decision_y)  # align horizontally with decision
            from_anchor = "left"
            to_anchor = "right"
            label_x = (decision_x + bx) / 2
            label_y = (decision_y + by) / 2 - 0.15  # CHANGED: closer to arrow midpoint
        elif branch_key == "path2":  # right box
            bx, by = (right_branch_x, decision_y)  # align horizontally with decision
            from_anchor = "right"
            to_anchor = "left"
            label_x = (decision_x + bx) / 2
            label_y = (decision_y + by) / 2 - 0.15  # CHANGED: closer to arrow midpoint
        elif branch_key == "path3":  # bottom box
            bx, by = (spine_x, bottom_action_y)
            from_anchor = "bottom"
            to_anchor = "top"
            label_x = spine_x + 0.3
            label_y = (decision_y + decision_h + by) / 2 - 0.15  # CHANGED: closer to arrow
        else:
            bx, by = (spine_x, decision_y + 1.0)
            from_anchor = "bottom"
            to_anchor = "top"
            label_x = bx
            label_y = by - 0.3
        
        lk = id(branch_root)
        if lk in memo:
            branch_id = memo[lk]
            w = node_geom[lk]["w"]
            h = node_geom[lk]["h"]
        else:
            branch_id = new_id(branch_root["type"].lower())
            memo[lk] = branch_id
            w, h = calculate_box_size(branch_root["text"], branch_root["type"])
            
            # CENTER branch boxes
            nodes.append({
                "id": branch_id,
                "type": branch_root["type"],
                "lines": branch_root["text"],
                "left": bx - w/2,  # CENTER
                "top": by,
                "width": w,
                "height": h,
            })
            node_geom[lk] = {"w": w, "h": h}
        
        connectors.append({
            "from_id": decision_id,
            "to_id": branch_id,
            "style": "straight",
            "from_anchor": from_anchor,
            "to_anchor": to_anchor,
            "label": label_text,
            "label_pos": (label_x, label_y),
        })
        
        cur = branch_root.get("paths", {}).get("next")
        prev_id = branch_id
        cur_y = by
        while cur:
            cur_y += y_step
            lk2 = id(cur)
            if lk2 in memo:
                nid2 = memo[lk2]
                w2 = node_geom[lk2]["w"]
                h2 = node_geom[lk2]["h"]
            else:
                nid2 = new_id(cur["type"].lower())
                memo[lk2] = nid2
                w2, h2 = calculate_box_size(cur["text"], cur["type"])
                
                # CENTER in branch column
                nodes.append({
                    "id": nid2,
                    "type": cur["type"],
                    "lines": cur["text"],
                    "left": bx - w2/2,  # CENTER
                    "top": cur_y,
                    "width": w2,
                    "height": h2,
                })
                node_geom[lk2] = {"w": w2, "h": h2}
            
            connectors.append({
                "from_id": prev_id,
                "to_id": nid2,
                "style": "straight",
                "from_anchor": "bottom",
                "to_anchor": "top",
            })
            
            prev_id = nid2
            cur = cur.get("paths", {}).get("next")
        
        return prev_id
    
    # 6. Place each branch
    branch_last_ids = {}
    for k in ["path1", "path2", "path3"]:
        if k in decision_paths:
            branch_last_ids[k] = place_branch(k)
    
    # 7. Every branch connects to END
    if final_end_node_id:
        for branch_key, branch_id in branch_last_ids.items():
            if branch_key == "path1":  # left branch
                from_anchor = "bottom"
                to_anchor = "left"
                prefer_route = "v-first"
            elif branch_key == "path2":  # right branch
                from_anchor = "bottom"
                to_anchor = "right"
                prefer_route = "v-first"
            elif branch_key == "path3":  # bottom branch
                from_anchor = "bottom"
                to_anchor = "top"
                prefer_route = "v-first"
            else:
                from_anchor = "center"
                to_anchor = "center"
                prefer_route = "h-first"
            
            connectors.append({
                "from_id": branch_id,
                "to_id": final_end_node_id,
                "style": "straight",
                "from_anchor": from_anchor,
                "to_anchor": to_anchor,
                "prefer": prefer_route,
            })
            print(f"✓ Added arrow: {branch_id} → {final_end_node_id} ({branch_key})")
    
    print(f"\n📊 Summary:")
    print(f"  - Total nodes: {len(nodes)}")
    print(f"  - Total connectors: {len(connectors)}")
    print(f"  - Branch boxes: {len(branch_last_ids)}")
    print(f"  - End node ID: {final_end_node_id}\n")
    
    return nodes, connectors


# Usage: Replace the generate_nodes_connectors_from_flow call with:
# nodes, connectors = generate_nodes_connectors_from_flow_flexible(flow_spec)


# =========================================================
# STEP 4. Drawing helpers (shapes, text, arrows)
# =========================================================
def rgb(rgb_tuple):
    return RGBColor(*rgb_tuple)

def pp_align(align_str):
    if align_str == "left":
        return PP_ALIGN.LEFT
    if align_str == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.CENTER

def add_box_shape(slide, node, styles):
    style = styles[node["type"]]
    fill_rgb = style["fill_rgb"]
    shape_kind = style["shape"]
    bold_first = style["bold_first_line"]
    alignment_pref = node.get("align", style.get("align", "center"))

    left = Inches(node["left"])
    top = Inches(node["top"])
    width = Inches(node["width"])
    height = Inches(node["height"])

    if shape_kind == "diamond":
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.DIAMOND,
            left, top, width, height
        )
    else:
        shp = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left, top, width, height
        )

    # disable shadow if any theme applies it
    try:
        shp.shadow.inherit = False
        shp.shadow.visible = False
    except Exception:
        pass

    # fill + outline
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_rgb)
    # Stronger outline and pill corners for Start/End
    if node["type"] in ("Start", "End"):
        try:
            # Make rounded rectangle more like a pill (0..1)
            if shp.adjustments and len(shp.adjustments) > 0:
                shp.adjustments[0] = 0.75
        except Exception:
            pass
        shp.line.width = Pt(2)
    else:
        shp.line.width = Pt(1.25)
    shp.line.color.rgb = rgb(DARK_TEAL)

    # text
    tf = shp.text_frame
    tf.clear()
    for i, line in enumerate(node["lines"]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(*FLOW_TEXT_RGB)
        if i == 0 and bold_first:
            run.font.bold = True
        p.alignment = pp_align(alignment_pref)

    return shp


def shape_center(shp):
    cx = (shp.left + shp.width / 2.0) / Inches(1)
    cy = (shp.top + shp.height / 2.0) / Inches(1)
    return cx, cy

def shape_top(shp):
    cx, _ = shape_center(shp)
    y = shp.top / Inches(1)
    return cx, y

def shape_bottom(shp):
    cx, _ = shape_center(shp)
    y = (shp.top + shp.height) / Inches(1)
    return cx, y

def shape_left_mid(shp):
    x = shp.left / Inches(1)
    _, cy = shape_center(shp)
    return x, cy

def shape_right_mid(shp):
    x = (shp.left + shp.width) / Inches(1)
    _, cy = shape_center(shp)
    return x, cy

def get_anchor_xy(shp, anchor):
    if anchor == "top":
        return shape_top(shp)
    if anchor == "bottom":
        return shape_bottom(shp)
    if anchor == "left":
        return shape_left_mid(shp)
    if anchor == "right":
        return shape_right_mid(shp)
    return shape_center(shp)



def add_label(slide, text, x, y):
    # small textbox near decision branch labels
    box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(1.8),
        Inches(0.4)
    )
    try:
        box.shadow.inherit = False
        box.shadow.visible = False
    except Exception:
        pass
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER
    return box


# =========================================================
# STEP 5. Build slide (title, legend, flowchart, footer)
# =========================================================
def build_slide_from_nodes_connectors(
    prs,
    title_text,
    logo_path,
    footer_text,
    nodes,
    connectors,
):
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title box
    title_box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(0.4),
        Inches(4.5),
        Inches(0.6)
    )
    try:
        title_box.shadow.inherit = False
        title_box.shadow.visible = False
    except Exception:
        pass
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title_text
    r.font.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = RGBColor(*DOC_HEADING_RGB)

    # Logo (optional)
    if logo_path:
        pic = slide.shapes.add_picture(
            logo_path,
            Inches(9.5),
            Inches(0.2),
            height=Inches(1.4)
        )
        try:
            pic.shadow.inherit = False
            pic.shadow.visible = False
        except Exception:
            pass

    # Legend ("Key") box
    key_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.8),
        Inches(1.2),
        Inches(3.0),
        Inches(1.4)
    )
    try:
        key_box.shadow.inherit = False
        key_box.shadow.visible = False
    except Exception:
        pass
    key_box.fill.solid()
    key_box.fill.fore_color.rgb = RGBColor(*KEY_BOX_RGB)
    key_box.line.width = Pt(1.25)
    key_box.line.color.rgb = rgb(DARK_TEAL)

    key_tf = key_box.text_frame
    key_tf.clear()

    p0 = key_tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = "Key:"
    run0.font.bold = True
    run0.font.size = Pt(12)
    run0.font.color.rgb = RGBColor(*KEY_TEXT_RGB)

    p1 = key_tf.add_paragraph()
    
    run1 = p1.add_run()
    run1.text = "Acronym 1: Description"
    run1.font.size = Pt(11)
    run1.font.color.rgb = RGBColor(*KEY_TEXT_RGB)

    p2 = key_tf.add_paragraph()
    
    run2 = p2.add_run()
    run2.text = "Acronym 2: Description"
    run2.font.size = Pt(11)
    run2.font.color.rgb = RGBColor(*KEY_TEXT_RGB)

    # Draw nodes first
    id_to_shape = {}
    for node in nodes:
        shp = add_box_shape(slide, node, BOX_STYLES)
        id_to_shape[node["id"]] = shp

    # Draw connectors (arrows) and then labels
    for conn in connectors:
        from_shape = id_to_shape[conn["from_id"]]
        to_shape   = id_to_shape[conn["to_id"]]

        # get anchor coords for start/end in inches
        x1, y1 = get_anchor_xy(from_shape, conn.get("from_anchor", "center"))
        x2, y2 = get_anchor_xy(to_shape,   conn.get("to_anchor", "center"))

        # Get label text if it exists
        label_text = conn.get("label", None)

        # draw the arrow connector WITH label on the line
        add_arrow_connector(
            slide=slide,
            x1=Inches(x1),
            y1=Inches(y1),
            x2=Inches(x2),
            y2=Inches(y2),
            line_width_pt=1,
            rgb=ARROW_RGB,
            arrow_at='end',
            label_text=label_text,  # Pass the label to be drawn ON the arrow
            prefer=conn.get("prefer", "v-first")
        )
       

        # If this connector had a label (Decision path 1 / 2 / 3), place it
        '''
        if "label" in conn and "label_pos" in conn:
            lx, ly = conn["label_pos"]
            add_label(slide, conn["label"], lx, ly)
        '''

    # Footer left (page number)
    page_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(7.1),
        Inches(1.0),
        Inches(0.3)
    )
    try:
        page_box.shadow.inherit = False
        page_box.shadow.visible = False
    except Exception:
        pass
    tfp = page_box.text_frame
    tfp.clear()
    pnum = tfp.paragraphs[0]
    rnum = pnum.add_run()
    rnum.text = "1"
    rnum.font.size = Pt(8)
    rnum.font.color.rgb = RGBColor(128, 128, 128)
    pnum.alignment = PP_ALIGN.LEFT

    # Footer right (metadata)
    footer_box = slide.shapes.add_textbox(
        Inches(7.5),
        Inches(7.1),
        Inches(5.5),
        Inches(0.3)
    )
    try:
        footer_box.shadow.inherit = False
        footer_box.shadow.visible = False
    except Exception:
        pass
    tff = footer_box.text_frame
    tff.clear()
    pf = tff.paragraphs[0]
    rf = pf.add_run()
    rf.text = footer_text
    rf.font.size = Pt(8)
    rf.font.color.rgb = RGBColor(*FOOTER_RGB)
    pf.alignment = PP_ALIGN.RIGHT

    return slide


# =========================================================
# STEP 6. Main script
# =========================================================
if __name__ == "__main__":
    # 1. parse bullets
    human_flow = parse_bullets_to_human_flow(BULLET_TEXT)

    # 2. human_flow -> flow_spec
    flow_spec = to_flow_spec(human_flow)

    # 3. flow_spec -> nodes/connectors
    nodes, connectors = nodes, connectors = generate_nodes_connectors_from_flow_flexible(flow_spec)

    # 4. build ppt
    prs = Presentation()

    TITLE_TEXT = "Document Name Here"
    LOGO_PATH = "my_power_logo.png"  # e.g. "mypower_logo.png" or "" for none
    FOOTER_TEXT = "Document Name – dd/mm/yy - Author: Your Name - Version: Draft"

    build_slide_from_nodes_connectors(
        prs=prs,
        title_text=TITLE_TEXT,
        logo_path=LOGO_PATH,
        footer_text=FOOTER_TEXT,
        nodes=nodes,
        connectors=connectors,
    )

    prs.save("process_flow_generated.pptx")
    print("Saved process_flow_generated.pptx")