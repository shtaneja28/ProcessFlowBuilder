# -*- coding: utf-8 -*-
"""
Further refinements:
- Stronger edge-hug avoidance: connectors aren't allowed to run along a box side.
  We enforce a small clearance pad around every box and exclude only the
  source/destination boxes for the current edge.
- Smaller Start/End lozenges (e.g., 1.6" x 0.55") with 20pt text.
- Keep decision labels at arrow origins.
- Single-slide serpentine; auto-fit if many columns.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from collections import defaultdict, deque
from datetime import datetime
from math import ceil
import sys
import argparse
import re

# ---------- Palette ----------
def hex_to_rgb(hex_str):
    hex_str = hex_str.strip().lstrip('#')
    return tuple(int(hex_str[i:i+2], 16) for i in (0,2,4))

PALETTE = {
    "document_heading": hex_to_rgb("005077"),
    "sub_heading": hex_to_rgb("ED1165"),
    "key_box": hex_to_rgb("92D050"),
    "key_text": hex_to_rgb("000000"),
    "flow_text": hex_to_rgb("000000"),
    "start_lozenge": hex_to_rgb("92D050"),
    "info_box": hex_to_rgb("2FC9FF"),
    "action_box": hex_to_rgb("9DC3E6"),
    "decision_box": hex_to_rgb("EAB0FA"),
    "end_lozenge": hex_to_rgb("FFC000"),
    "arrows": hex_to_rgb("4472C4"),
    "footer": hex_to_rgb("BFBFBF"),
    "outline": hex_to_rgb("000000"),
}

# ---------- Footer helper ----------
def _add_standard_footer(prs, slide):
    """Adds left page count and right document info footer to a slide."""
    page_num = len(prs.slides)
    left_fb = slide.shapes.add_textbox(Inches(0.2), prs.slide_height - Inches(0.35), Inches(2.0), Inches(0.3))
    ltf = left_fb.text_frame; ltf.clear()
    lp = ltf.paragraphs[0]; lp.alignment = PP_ALIGN.LEFT
    lr = lp.add_run(); lr.text = f"Page {page_num} of xx"; lr.font.size = Pt(9); lr.font.color.rgb = RGBColor(*PALETTE["footer"])

    today_str = datetime.now().strftime("%d/%m/%y")
    right_text = f"PROC No  Document Name – {today_str} - Author: Your Name - Version: Draft"
    rb = slide.shapes.add_textbox(Inches(0.3), prs.slide_height - Inches(0.45), prs.slide_width - Inches(0.6), Inches(0.3))
    rtf = rb.text_frame; rtf.clear()
    rp = rtf.paragraphs[0]; rp.alignment = PP_ALIGN.RIGHT
    rr = rp.add_run(); rr.text = right_text; rr.font.size = Pt(9); rr.font.color.rgb = RGBColor(*PALETTE["footer"])

# ---------- Key box helper ----------
def _add_key_box(prs, slide, key_lines, logo_present=False):
    """Add a 'Key' box to the top-right of the flow slide.
    key_lines: list[str] of lines to display under the Key title.
    """
    box_w = Inches(3.0)
    left = prs.slide_width - Inches(0.6) - box_w
    # If there is a logo at the top-right at ~0.2..1.2in tall, drop the key box slightly
    top = Inches(1.3) if logo_present else Inches(0.95)

    # Dynamic height estimation using PIL (96 DPI assumptions)
    def _load_font_local(pt):
        candidates = [
            "/System/Library/Fonts/Supplemental/Arial.ttf",
            "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Supplemental/Helvetica.ttc",
            "/System/Library/Fonts/Supplemental/Calibri.ttf",
            "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/calibri.ttf",
        ]
        for p in candidates:
            try:
                return ImageFont.truetype(p, pt)
            except Exception:
                continue
        return ImageFont.load_default()

    def _wrap_lines_local(text, font, max_px):
        tokens = []
        for raw in text.split():
            tmp = raw.replace('/', ' / ').replace('-', ' - ')
            tokens.extend(tmp.split())
        if not tokens:
            return [""]
        lines_out = []
        cur = tokens[0]
        for t in tokens[1:]:
            trial = f"{cur} {t}"
            if font.getlength(trial) <= max_px:
                cur = trial
            else:
                if font.getlength(t) > max_px:
                    if cur:
                        lines_out.append(cur)
                        cur = ""
                    piece = ""
                    for ch in t:
                        nxt = piece + ch
                        if font.getlength(nxt) <= max_px:
                            piece = nxt
                        else:
                            if piece:
                                lines_out.append(piece)
                            piece = ch
                    cur = piece if piece else ""
                else:
                    lines_out.append(cur)
                    cur = t
        if cur:
            lines_out.append(cur)
        return lines_out

    dpi = 96.0
    inner_w_in = float(box_w - Inches(0.2))
    max_px = int(inner_w_in * dpi)
    font_title = _load_font_local(12)
    font_body = _load_font_local(9)
    asc, dsc = font_body.getmetrics()
    line_h_px = asc + dsc + 2

    total_lines = 1  # title 'Key'
    for raw in key_lines:
        wrapped = _wrap_lines_local((raw or "").strip(), font_body, max_px)
        total_lines += max(1, len(wrapped))
    content_h_in = (total_lines * line_h_px) / dpi
    # margins + small safety
    box_h = Inches(0.28) + Inches(content_h_in * 1.05)
    min_h = Inches(1.2)
    max_h = prs.slide_height - top - Inches(0.9)
    if box_h < min_h:
        box_h = min_h
    if box_h > max_h:
        box_h = max_h

    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, box_w, box_h)
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(*PALETTE["key_box"])  # light highlight
    shp.line.color.rgb = RGBColor(*PALETTE["outline"])
    shp.line.width = Pt(1.25)
    try:
        shp.shadow.visible = False
        shp.shadow.inherit = False
    except Exception:
        pass

    tf = shp.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.TOP
    # Title
    p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run(); r0.text = "Key"; r0.font.size = Pt(12); r0.font.bold = True; r0.font.underline = False; r0.font.color.rgb = RGBColor(*PALETTE["key_text"])
    # Body lines
    for line in key_lines:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        txt = (line or "").rstrip("\n")
        if not txt:
            p.text = ""  # blank line
            continue
        if txt.strip().startswith("•"):
            p.level = 0; p.text = txt.strip()
            for run in p.runs:
                run.font.size = Pt(9); run.font.color.rgb = RGBColor(*PALETTE["key_text"])
        else:
            r = p.add_run(); r.text = txt; r.font.size = Pt(9); r.font.color.rgb = RGBColor(*PALETTE["key_text"])
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.06)

# ---------- Parse explicit-ID schema ----------
NODE_RE = re.compile(r'^(Start|End|Decision|Action|Info|Information):\s*\[(.+?)\]\s*(.*)$', re.I)
DETAILS_RE = re.compile(r'^\s*Details:\s*(.*)$', re.I)
TITLE_RE = re.compile(r'^\s*Title:\s*(.*)$', re.I)
LEADS_RE = re.compile(r'^\s*Leads to:\s*\[(.+?)\]\s*$', re.I)
PATH_RE = re.compile(r'^\s*Path\s+"(.+?)"\s*->\s*\[(.+?)\]\s*$', re.I)

def parse_schema_details_only(text):
    nodes = {}
    edges = []
    start_ids = []
    current_id = None
    current_kind = None

    lines = [l.rstrip("\n") for l in text.splitlines()]
    for line in lines:
        if not line.strip():
            continue

        m = NODE_RE.match(line)
        if m:
            kind = m.group(1).lower()
            if kind == "information":
                kind = "info"
            nid = m.group(2).strip()
            rest = m.group(3).strip()
            current_id = nid
            current_kind = kind
            if nid not in nodes:
                nodes[nid] = {"id": nid, "kind": kind, "title": "", "details_lines": []}
            if kind in ("start","end","decision") and rest:
                nodes[nid]["title"] = rest
            if kind == "start":
                start_ids.append(nid)
            continue

        m = DETAILS_RE.match(line)
        if m and current_id:
            nodes[current_id]["details_lines"].append(m.group(1))
            continue

        m = TITLE_RE.match(line)
        if m and current_id:
            nodes[current_id]["title"] = m.group(1).strip()
            continue

        m = LEADS_RE.match(line)
        if m and current_id:
            edges.append((current_id, m.group(1).strip(), ""))
            continue

        m = PATH_RE.match(line)
        if m and current_id and current_kind == "decision":
            edges.append((current_id, m.group(2).strip(), m.group(1).strip()))
            continue

    return nodes, edges, start_ids

# ---------- Decision routing ----------
def plan_decision_routes(nodes, edges):
    out_by_src = defaultdict(list)
    for u, v, lbl in edges:
        out_by_src[u].append((v, lbl))

    route_pref = {}
    for nid, ndata in nodes.items():
        if ndata["kind"] != "decision":
            continue
        outs = out_by_src.get(nid, [])
        if not outs: 
            continue
        outs = sorted(outs, key=lambda t: (t[1] != "Yes", t[1]))
        if len(outs) == 1:
            route_pref[(nid, outs[0][0])] = "down"
        else:
            route_pref[(nid, outs[0][0])] = "right"
            route_pref[(nid, outs[1][0])] = "down"
    return route_pref

# ---------- Columns & order ----------
from collections import deque, defaultdict

def assign_columns(nodes, edges, start_ids, route_pref):
    if not nodes:
        return {}
    col = {nid: None for nid in nodes}
    q = deque()
    for s in start_ids or [next(iter(nodes))]:
        col[s] = 0; q.append(s)
    adj = defaultdict(list); indeg = defaultdict(int)
    for u, v, _ in edges:
        adj[u].append(v); indeg[v] += 1
    while q:
        u = q.popleft()
        cu = col[u]
        for v in adj.get(u, []):
            pref = route_pref.get((u, v))
            dv = cu + 1 if (nodes[u]["kind"] == "decision" and pref == "right") else cu
            if col[v] is None or dv > col[v]:
                col[v] = dv
            indeg[v] -= 1
            if indeg[v] <= 0:
                q.append(v)
    base = min(c for c in col.values() if c is not None)
    for k in col:
        col[k] = (col[k] or base) - base
    return col

def order_within_columns(nodes, edges, columns):
    out = defaultdict(list); indeg = defaultdict(int)
    for u, v, _ in edges:
        out[u].append(v); indeg[v] += 1
    q = deque([n for n in nodes if indeg[n] == 0])
    topo = []
    seen = set()
    while q:
        u = q.popleft()
        if u in seen: continue
        seen.add(u); topo.append(u)
        for v in out.get(u, []):
            indeg[v] -= 1
            if indeg[v] <= 0:
                q.append(v)
    by_col = defaultdict(list); rank = {nid: i for i, nid in enumerate(topo)}
    for nid in nodes:
        by_col[columns[nid]].append(nid)
    for c in by_col:
        by_col[c].sort(key=lambda n: rank.get(n, 0))
    return by_col

# ---------- Draw helpers ----------
from pptx.oxml.xmlchemy import OxmlElement
from PIL import ImageFont

def add_node_shape(slide, node, left, top, width, height, start_end_override=None):
    kind = node["kind"]
    title = node["title"]
    details = node["details_lines"]

    if start_end_override and kind in ("start","end"):
        sw, sh = start_end_override
        left = left + (width - sw)//2
        width = sw
        height = sh

    if kind == "decision":
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.DIAMOND, left, top, width, max(height, Inches(1.0)))
        fill = PALETTE["decision_box"]; align = PP_ALIGN.CENTER; bold = False; fsize = 9
    elif kind == "start":
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        try: shp.adjustments[0] = 0.7
        except Exception: pass
        fill = PALETTE["start_lozenge"]; align = PP_ALIGN.CENTER; bold = True; fsize = 20
    elif kind == "end":
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        try: shp.adjustments[0] = 0.7
        except Exception: pass
        fill = PALETTE["end_lozenge"]; align = PP_ALIGN.CENTER; bold = True; fsize = 20
    elif kind == "info":
        # Information tag: rounded rectangle with info_box colour
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
        try: shp.adjustments[0] = 0.1
        except Exception: pass
        fill = PALETTE["info_box"]; align = PP_ALIGN.LEFT; bold = False; fsize = 9
    else:
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
        fill = PALETTE["action_box"]; align = PP_ALIGN.LEFT; bold = False; fsize = 9

    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(*fill)
    shp.line.color.rgb = RGBColor(*PALETTE["outline"]); shp.line.width = Pt(1.0 if kind!="decision" else 1.25)
    # Ensure no shadows on shapes
    try:
        shp.shadow.visible = False
        shp.shadow.inherit = False
    except Exception:
        pass

    tf = shp.text_frame; tf.clear()
    if kind in ("decision","start","end"):
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = title; r.font.size = Pt(fsize); r.font.color.rgb = RGBColor(*PALETTE["flow_text"]); r.font.bold = bold
    else:
        first = True
        # Optional Title heading for action boxes
        if title:
            p = tf.paragraphs[0]; first = False
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = title
            r.font.size = Pt(9); r.font.color.rgb = RGBColor(*PALETTE["flow_text"]); r.font.bold = True; r.font.underline = True
        for line in details:
            if first:
                p = tf.paragraphs[0]; first = False
            else:
                p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            if line.strip().startswith("•"):
                p.level = 0; p.text = line.strip()
                for run in p.runs:
                    run.font.size = Pt(9); run.font.color.rgb = RGBColor(*PALETTE["flow_text"])
            else:
                r = p.add_run(); r.text = line
                r.font.size = Pt(9); r.font.color.rgb = RGBColor(*PALETTE["flow_text"])
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.1)
    tf.margin_top = tf.margin_bottom = Inches(0.06)
    # Dynamically fit text for content-heavy boxes while keeping Start/End fixed at 20pt
    if kind not in ("start", "end"):
        try:
            tf.fit_text(max_size=9)
        except Exception:
            pass
        # Ensure Title (when present) remains bold and exactly 9pt after any fit_text scaling
        if title and kind not in ("decision",):
            try:
                p0 = tf.paragraphs[0]
                for run in p0.runs:
                    run.font.size = Pt(9)
                    run.font.bold = True
                    run.font.underline = True
            except Exception:
                pass
    return shp

def mid_top(s):    return (s.left + s.width//2, s.top)
def mid_bottom(s): return (s.left + s.width//2, s.top + s.height)
def mid_left(s):   return (s.left, s.top + s.height//2)
def mid_right(s):  return (s.left + s.width, s.top + s.height//2)

def add_seg(slide, x1,y1,x2,y2, rgb, end=False):
    kind = MSO_CONNECTOR.STRAIGHT if (x1==x2 or y1==y2) else MSO_CONNECTOR.ELBOW
    c = slide.shapes.add_connector(kind, x1,y1,x2,y2)
    c.line.fill.solid(); c.line.fill.fore_color.rgb = RGBColor(*rgb); c.line.width = Pt(2.0)
    ln = c.line._get_or_add_ln()
    if end:
        e = OxmlElement('a:tailEnd'); e.set('type','arrow'); e.set('w','sm'); e.set('len','sm'); ln.append(e)
    # Disable shadows on connectors as well
    try:
        c.shadow.visible = False
        c.shadow.inherit = False
    except Exception:
        pass
    return c

def add_seg_registered(slide, x1,y1,x2,y2, rgb, end, used_h, used_v):
    add_seg(slide, x1, y1, x2, y2, rgb, end=end)
    if x1 == x2:
        y_low, y_high = sorted([y1, y2])
        used_v.add((x1, y_low, y_high))
    elif y1 == y2:
        x_low, x_high = sorted([x1, x2])
        used_h.add((y1, x_low, x_high))
    return True

def segment_hits_rect(x1,y1,x2,y2, rect, pad):
    nx, ny, nw, nh = rect
    nx -= pad; ny -= pad; nw += 2*pad; nh += 2*pad
    if x1 == x2:
        top, bottom = sorted([y1,y2]); x = x1
        # <= and >= so running along the side counts as a hit
        return (nx <= x <= nx+nw) and not (ny+nh <= top or ny >= bottom)
    elif y1 == y2:
        left, right = sorted([x1,x2]); y = y1
        return (ny <= y <= ny+nh) and not (nx+nw <= left or nx >= right)
    else:
        return True

def path_clear(a,b, rects, pad):
    for r in rects:
        if segment_hits_rect(a[0],a[1], b[0],b[1], r, pad):
            return False
    return True

def route_orthogonal_detour(slide, p1, p2, rects_all, rgb, exclude_rects=(), used_h=None, used_v=None, ignore_lanes=False):
    """Try L-shape, else gutters; avoid running along box sides and existing lanes.
    If ignore_lanes is True, skip lane collision checks (useful for shared/merged routes).
    """
    PAD = Inches(0.05)  # ~0.05" clearance around every box
    rects = [r for r in rects_all if r not in exclude_rects]

    elbow1 = (p2[0], p1[1])  # H then V
    elbow2 = (p1[0], p2[1])  # V then H
    used_h = used_h if used_h is not None else set()
    used_v = used_v if used_v is not None else set()

    def lane_free(a, b):
        if ignore_lanes:
            return True
        x1,y1 = a; x2,y2 = b
        if x1 == x2:
            y_low, y_high = sorted([y1, y2])
            for (x, yl, yh) in used_v:
                if x == x1 and not (yh <= y_low or yl >= y_high):
                    return False
        elif y1 == y2:
            x_low, x_high = sorted([x1, x2])
            for (y, xl, xh) in used_h:
                if y == y1 and not (xh <= x_low or xl >= x_high):
                    return False
        return True

    if path_clear(p1, elbow1, rects, PAD) and path_clear(elbow1, p2, rects, PAD) and lane_free(p1, elbow1) and lane_free(elbow1, p2):
        add_seg_registered(slide, p1[0],p1[1], elbow1[0],elbow1[1], rgb, end=False, used_h=used_h, used_v=used_v)
        add_seg_registered(slide, elbow1[0],elbow1[1], p2[0],p2[1], rgb, end=True, used_h=used_h, used_v=used_v)
        return True
    if path_clear(p1, elbow2, rects, PAD) and path_clear(elbow2, p2, rects, PAD) and lane_free(p1, elbow2) and lane_free(elbow2, p2):
        add_seg_registered(slide, p1[0],p1[1], elbow2[0],elbow2[1], rgb, end=False, used_h=used_h, used_v=used_v)
        add_seg_registered(slide, elbow2[0],elbow2[1], p2[0],p2[1], rgb, end=True, used_h=used_h, used_v=used_v)
        return True

    # Gutters lines between boxes
    xs = sorted(set([nx - Inches(0.12) for (nx,ny,nw,nh) in rects] + [nx+nw + Inches(0.12) for (nx,ny,nw,nh) in rects]))
    ys = sorted(set([ny - Inches(0.12) for (nx,ny,nw,nh) in rects] + [ny+nh + Inches(0.12) for (nx,ny,nw,nh) in rects]))

    for yg in ys:
        a = (p1[0], p1[1]); b = (p1[0], yg); c = (p2[0], yg); d = (p2[0], p2[1])
        if path_clear(a, b, rects, PAD) and path_clear(b, c, rects, PAD) and path_clear(c, d, rects, PAD) and lane_free(a, b) and lane_free(b, c) and lane_free(c, d):
            add_seg_registered(slide, a[0],a[1], b[0],b[1], rgb, end=False, used_h=used_h, used_v=used_v)
            add_seg_registered(slide, b[0],b[1], c[0],c[1], rgb, end=False, used_h=used_h, used_v=used_v)
            add_seg_registered(slide, c[0],c[1], d[0],d[1], rgb, end=True, used_h=used_h, used_v=used_v)
            return True

    for xg in xs:
        a = (p1[0], p1[1]); b = (xg, p1[1]); c = (xg, p2[1]); d = (p2[0], p2[1])
        if path_clear(a, b, rects, PAD) and path_clear(b, c, rects, PAD) and path_clear(c, d, rects, PAD) and lane_free(a, b) and lane_free(b, c) and lane_free(c, d):
            add_seg_registered(slide, a[0],a[1], b[0],b[1], rgb, end=False, used_h=used_h, used_v=used_v)
            add_seg_registered(slide, b[0],b[1], c[0],c[1], rgb, end=False, used_h=used_h, used_v=used_v)
            add_seg_registered(slide, c[0],c[1], d[0],d[1], rgb, end=True, used_h=used_h, used_v=used_v)
            return True

    # last resort
    add_seg_registered(slide, p1[0],p1[1], p2[0],p2[1], rgb, end=True, used_h=used_h, used_v=used_v)
    return False

# ---------- Render ----------
def render(schema_text, output_path, logo_path=None, show_key=False, key_path=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    # Insert cover slide first
    _add_cover_slide(prs, logo_path)
    # Insert amendments slide second
    _add_amendments_slide(prs, logo_path)
    # Flow slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header + logo (logo top-right, title on the left)
    title_text = "Document Name Here"
    title_top = Inches(0.3); title_left = Inches(0.6); title_height = Inches(0.6)
    title_width = prs.slide_width - Inches(1.2)
    logo_present = False
    if logo_path:
        logo_height = Inches(1.0)
        slide.shapes.add_picture(str(logo_path), prs.slide_width - Inches(0.6) - logo_height*2.0, Inches(0.2), height=logo_height)
        title_width = prs.slide_width - Inches(1.2) - (logo_height*2.0 + Inches(0.2))
        logo_present = True
    tb = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title_text; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RGBColor(*PALETTE["document_heading"])

    # Key box (optional)
    if show_key and key_path:
        try:
            from pathlib import Path as _P
            kp = _P(key_path)
            if kp.exists():
                with open(kp, "r", encoding="utf-8") as kf:
                    key_lines = kf.readlines()
                _add_key_box(prs, slide, key_lines, logo_present=logo_present)
        except Exception:
            # Fail-soft: do not block rendering if key cannot be read
            pass

    # Parse, route prefs, columns
    nodes, edges, start_ids = parse_schema_details_only(schema_text)
    route_pref = plan_decision_routes(nodes, edges)
    columns = assign_columns(nodes, edges, start_ids, route_pref)
    if not columns:
        # No nodes parsed; print a helpful note and return early
        print("No nodes parsed from schema – check the schema file format and path.")
        prs.save(output_path)
        return output_path
    by_col = order_within_columns(nodes.keys(), edges, columns)
    base_col_count = max(columns.values()) + 1

    # Sizing
    cfg = {
        "margin_l": Inches(0.6),
        "margin_r": Inches(0.6),
        "margin_t": Inches(1.1),
        "margin_b": Inches(0.6),
        "v_gap": Inches(0.3),
        "h_gap": Inches(1.0),
        "box_w": Inches(3.0),
        "box_h": Inches(0.9),  # legacy default; dynamic per action below
        "box_h_min": Inches(0.7),
        "line_height_in": 0.20,  # slightly larger to reduce overflow risk
        "start_end_w": Inches(1.6),
        "start_end_h": Inches(0.55),
    }
    usable_h = prs.slide_height - cfg["margin_t"] - cfg["margin_b"]
    usable_w = prs.slide_width  - cfg["margin_l"] - cfg["margin_r"]

    # Serpentine split for tall columns
    fixed_order_cols = {}
    total_columns = 0
    for c in range(base_col_count):
        items = by_col.get(c, [])
        if not items: continue
        per_col = max(1, int((usable_h + cfg["v_gap"]) // (cfg["box_h"] + cfg["v_gap"])))
        chunks = [items[i:i+per_col] for i in range(0, len(items), per_col)]
        fixed_order_cols[c] = chunks
        total_columns += len(chunks)

    # Auto scale horizontally
    needed_w = total_columns*cfg["box_w"] + (total_columns-1)*cfg["h_gap"]
    if needed_w > usable_w:
        scale = float(usable_w) / float(needed_w)
        cfg["box_w"] = int(cfg["box_w"] * scale)
        cfg["h_gap"] = int(cfg["h_gap"] * scale)
        cfg["start_end_w"] = int(cfg["start_end_w"] * scale)

    # Identify decisions that route downward to add extra vertical clearance for arrows/labels
    decisions_with_down = set()
    for (u, v, _lbl) in edges:
        if nodes[u]["kind"] == "decision" and route_pref.get((u, v)) == "down":
            decisions_with_down.add(u)

    def _load_font(pt: int) -> ImageFont.FreeTypeFont:
        # Try common fonts; fall back to a default bitmap font (approx size)
        candidates = [
            "/System/Library/Fonts/Supplemental/Arial.ttf",  # macOS
            "/Library/Fonts/Arial.ttf",
            "/System/Library/Fonts/Supplemental/Helvetica.ttc",
            "/System/Library/Fonts/Supplemental/Calibri.ttf",
            "C:/Windows/Fonts/arial.ttf",
            "C:/Windows/Fonts/calibri.ttf",
        ]
        for path in candidates:
            try:
                return ImageFont.truetype(path, pt)
            except Exception:
                continue
        return ImageFont.load_default()

    def _wrap_lines(text: str, font: ImageFont.ImageFont, max_px: int):
        # Split on whitespace and common separators to improve wrap quality
        tokens = []
        for raw in text.split():
            # Further split on '/' and '-' so very long tokens can break naturally
            parts = []
            tmp = raw.replace('/', ' / ').replace('-', ' - ')
            for p in tmp.split():
                parts.append(p)
            tokens.extend(parts)

        if not tokens:
            return [""]

        lines: list[str] = []
        cur = tokens[0]
        for t in tokens[1:]:
            trial = f"{cur} {t}"
            if font.getlength(trial) <= max_px:
                cur = trial
            else:
                # If a single token is too wide, break it at character level
                if font.getlength(t) > max_px:
                    # flush current line
                    if cur:
                        lines.append(cur)
                        cur = ""
                    piece = ""
                    for ch in t:
                        nxt = piece + ch
                        if font.getlength(nxt) <= max_px:
                            piece = nxt
                        else:
                            if piece:
                                lines.append(piece)
                            piece = ch
                    cur = piece if piece else ""
                else:
                    lines.append(cur)
                    cur = t
        if cur:
            lines.append(cur)
        return lines

    def estimate_action_height(nid):
        details = nodes[nid]["details_lines"]
        has_title = bool(nodes[nid].get("title"))
        # inside width in pixels at 96 DPI
        dpi = 96.0
        inner_w_in = float(cfg["box_w"] - Inches(0.2))
        max_px = int(inner_w_in * dpi)
        font_body = _load_font(9)
        font_title = _load_font(9)
        ascent, descent = font_body.getmetrics()
        line_h_px = ascent + descent + 2  # small padding per line

        total_lines = 0
        if has_title:
            total_lines += len(_wrap_lines(nodes[nid]["title"], font_title, max_px))
        for line in details:
            txt = (line or "").strip()
            wrapped = _wrap_lines(txt, font_body, max_px)
            total_lines += max(1, len(wrapped))

        # margins ~0.14" top/bottom
        content_h_in = (total_lines * line_h_px) / dpi
        # Add a safety factor so PPT's layout differences don't cause overflow
        h_in = 0.32 + content_h_in * 1.18
        return max(cfg["box_h_min"], Inches(h_in))

    def node_height(nid):
        kind = nodes[nid]["kind"]
        if kind == "decision":
            return max(cfg["box_h"], Inches(1.1))
        if kind in ("start", "end"):
            return cfg["start_end_h"]
        # action/info boxes: dynamic height based on content
        return estimate_action_height(nid)

    # Place nodes
    id_to_shape = {}
    x = cfg["margin_l"]
    for c in range(base_col_count):
        chunks = fixed_order_cols.get(c, [])
        for chunk in chunks:
            rows = len(chunk)
            # Compute total height taking into account decision height and extra gap after downward decisions
            extra_after_decision = Inches(0.25)
            total_h = 0
            for i, nid in enumerate(chunk):
                total_h += node_height(nid)
                if i < rows - 1:
                    total_h += cfg["v_gap"]
                    if nid in decisions_with_down:
                        total_h += extra_after_decision
            y0 = cfg["margin_t"] + max(0, (usable_h - total_h)//2)
            y = y0
            for nid in chunk:
                node = nodes[nid]
                h = node_height(nid)
                start_end_override = (cfg["start_end_w"], cfg["start_end_h"]) if node["kind"] in ("start","end") else None
                shp = add_node_shape(slide, node, x, y, cfg["box_w"], h, start_end_override=start_end_override)
                id_to_shape[nid] = shp
                # Increment with actual height and gaps; add extra clearance after downward decisions
                y += h + cfg["v_gap"]
                if nid in decisions_with_down:
                    y += extra_after_decision
            x += cfg["box_w"] + cfg["h_gap"]

    # Rects with ids for collision tests
    rects_by_id = { nid: (sh.left, sh.top, sh.width, sh.height) for nid, sh in id_to_shape.items() }
    rects_all = list(rects_by_id.values())

    # Draw connectors with start-side labels, edge-hug avoidance, and lane reservation (no intersections)
    used_h = set()  # (y, x1, x2)
    used_v = set()  # (x, y1, y2)
    incoming_counts = defaultdict(int)
    for (_uu, vv, _ll) in edges:
        incoming_counts[vv] += 1
    dest_entry_lane = {}  # (dest_id, side) -> (attach_point, lane_point)
    for (u,v,label) in edges:
        su, sv = id_to_shape[u], id_to_shape[v]
        pref = route_pref.get((u, v))
        if pref == "right":
            # Approach destination from the left using a shared horizontal lane so
            # multiple connectors to the same destination coincide.
            p1 = (su.left + su.width, su.top + su.height//2)
            attach_point = (sv.left, sv.top + sv.height//2)
            key = (v, "left")
            if key not in dest_entry_lane:
                lane_point = (sv.left - Inches(0.6), sv.top + sv.height//2)
                dest_entry_lane[key] = (attach_point, lane_point)
            attach_point, shared_lane = dest_entry_lane[key]
            p2 = attach_point
            label_offset = (Inches(0.1), -Inches(0.25))
        elif pref == "down":
            p1, p2 = (su.left + su.width//2, su.top + su.height), (sv.left + sv.width//2, sv.top)
            label_offset = (-Inches(0.35), Inches(0.05))
        else:
            if (sv.left > su.left + su.width):
                p1 = (su.left + su.width, su.top + su.height//2)
                attach_point = (sv.left, sv.top + sv.height//2)
                key = (v, "left")
                if key not in dest_entry_lane:
                    lane_point = (sv.left - Inches(0.6), sv.top + sv.height//2)
                    dest_entry_lane[key] = (attach_point, lane_point)
                attach_point, shared_lane = dest_entry_lane[key]
                p2 = attach_point
                label_offset = (Inches(0.1), -Inches(0.25))
            elif (sv.top > su.top):
                p1, p2 = (su.left + su.width//2, su.top + su.height), (sv.left + sv.width//2, sv.top)
                label_offset = (-Inches(0.35), Inches(0.05))
            else:
                p1 = (su.left + su.width, su.top + su.height//2)
                attach_point = (sv.left, sv.top + sv.height//2)
                key = (v, "left")
                if key not in dest_entry_lane:
                    lane_point = (sv.left - Inches(0.6), sv.top + sv.height//2)
                    dest_entry_lane[key] = (attach_point, lane_point)
                attach_point, shared_lane = dest_entry_lane[key]
                p2 = attach_point
                label_offset = (Inches(0.1), -Inches(0.25))

        # Exclude source & target rects when testing for collisions so we can depart/arrive on edges
        exclude = (rects_by_id[u], rects_by_id[v])
        allow_share = incoming_counts[v] > 1
        if 'shared_lane' in locals():
            route_orthogonal_detour(
                slide, p1, shared_lane, rects_all, PALETTE["arrows"],
                exclude_rects=exclude, used_h=used_h, used_v=used_v,
                ignore_lanes=allow_share
            )
            add_seg_registered(slide, shared_lane[0], shared_lane[1], p2[0], p2[1], PALETTE["arrows"], end=True, used_h=used_h, used_v=used_v)
            del shared_lane
        else:
            route_orthogonal_detour(slide, p1, p2, rects_all, PALETTE["arrows"], exclude_rects=exclude, used_h=used_h, used_v=used_v)

        if nodes[u]["kind"] == "decision" and label:
            # Place labels offset perpendicular to the emerging line so they don't cloud the arrow.
            if pref == "right" or (pref is None and (sv.left > su.left + su.width)):
                # Right path: keep label very close to the box and higher above the line
                lbl_off = (-Inches(0.10), -Inches(0.30))
            else:  # down or fallback
                # Down path: move label slightly left and below the line
                lbl_off = (-Inches(0.15), Inches(0.14))

            lx = p1[0] + lbl_off[0]
            ly = p1[1] + lbl_off[1]
            lab = slide.shapes.add_textbox(lx, ly, Inches(0.6), Inches(0.28))
            lab.fill.solid(); lab.fill.fore_color.rgb = RGBColor(255, 255, 255)
            try:
                lab.line.fill.background()
            except Exception:
                pass
            tf2 = lab.text_frame; tf2.clear()
            pl = tf2.paragraphs[0]; pl.alignment = PP_ALIGN.CENTER
            rn = pl.add_run(); rn.text = label; rn.font.size = Pt(9); rn.font.color.rgb = RGBColor(*PALETTE["flow_text"])

    # Footer
    _add_standard_footer(prs, slide)

    # Final notes slide matching the provided template/image
    _add_notes_slide(prs, logo_path)

    prs.save(output_path)
    return output_path

# ---------- Cover slide ----------
def _add_cover_slide(prs, logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Logo top-right and title on the left
    title_top = Inches(0.4)
    title_left = Inches(0.6)
    title_width = prs.slide_width - Inches(1.2)
    if logo_path:
        try:
            logo_h = Inches(1.0)
            slide.shapes.add_picture(str(logo_path), prs.slide_width - Inches(0.6) - logo_h*2.0, Inches(0.2), height=logo_h)
            title_width = prs.slide_width - Inches(1.2) - (logo_h*2.0 + Inches(0.2))
        except Exception:
            pass
    tb = slide.shapes.add_textbox(title_left, title_top, title_width, Inches(0.8))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Document Name Here"; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RGBColor(*PALETTE["document_heading"])

    # Version table (4 rows x 2 cols) - lifted up
    tbl_left = Inches(3.0); tbl_top = Inches(1.0); tbl_w = Inches(3.8); tbl_h = Inches(1.4)
    rows, cols = 4, 2
    table = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
    labels = ["Version No", "Author", "Date Approved", "Date Distributed"]
    for i in range(rows):
        cell_label = table.cell(i,0)
        cell_val = table.cell(i,1)
        # Label style
        cell_label.fill.solid(); cell_label.fill.fore_color.rgb = RGBColor(*hex_to_rgb("ED1165"))
        cell_label.text_frame.clear(); pl = cell_label.text_frame.paragraphs[0]
        rl = pl.add_run(); rl.text = labels[i]; rl.font.bold = True; rl.font.size = Pt(9); rl.font.color.rgb = RGBColor(0,0,0)
        # Value style
        cell_val.fill.solid(); cell_val.fill.fore_color.rgb = RGBColor(*hex_to_rgb("FBE6E1"))
        cell_val.text_frame.clear(); pv = cell_val.text_frame.paragraphs[0]
        rv = pv.add_run(); rv.text = "Draft" if i==0 else ""
        rv.font.size = Pt(9); rv.font.bold = True if i==0 else False; rv.font.color.rgb = RGBColor(0,0,0)
        # Borders
        for c in (cell_label, cell_val):
            tc = c._tc.get_or_add_tcPr()
            for tag in ('a:lnL','a:lnR','a:lnT','a:lnB'):
                ln = OxmlElement(tag); ln.set('w','9525')
                sf = OxmlElement('a:solidFill'); clr = OxmlElement('a:srgbClr'); clr.set('val','000000'); sf.append(clr); ln.append(sf)
                tc.append(ln)
    # Footer
    _add_standard_footer(prs, slide)
    # Distribution list title bar across content width
    content_left = Inches(0.6)
    content_right = prs.slide_width - Inches(0.6)
    gap_after_version = Inches(0.2)
    available_w = content_right - content_left
    desired_w = Inches(9.0)
    bar_w = desired_w if desired_w < available_w else available_w
    bar_left = content_left + (available_w - bar_w)//2
    bar_top = tbl_top + tbl_h + gap_after_version
    bar_h = Inches(0.35)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, bar_left, bar_top, bar_w, bar_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb("ED1165"))
    bar.line.color.rgb = RGBColor(0,0,0)
    t = bar.text_frame; t.clear(); pp = t.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = "Distribution List"; rr.font.bold = True; rr.font.size = Pt(9); rr.font.color.rgb = RGBColor(0,0,0)

    # Distribution table (headers + several rows)
    d_rows, d_cols = 15, 4
    dt_left = bar_left; dt_top = bar_top + bar_h
    dt_w = bar_w
    # Fit table height to slide with bottom margin
    dt_h = prs.slide_height - dt_top - Inches(0.6)
    dtable = slide.shapes.add_table(d_rows, d_cols, dt_left, dt_top, dt_w, dt_h).table
    headers = ["Tick to\nInclude", "Name", "Position", "Department"]
    for j in range(d_cols):
        hc = dtable.cell(0,j)
        hc.fill.solid(); hc.fill.fore_color.rgb = RGBColor(*hex_to_rgb("ED1165"))
        hc.text_frame.clear(); hp = hc.text_frame.paragraphs[0]; hr = hp.add_run(); hr.text = headers[j]; hr.font.bold = True; hr.font.size = Pt(9); hr.font.color.rgb = RGBColor(0,0,0)
        hp.alignment = PP_ALIGN.CENTER
    # Hardcoded distribution list rows (✔, Name, Position, Department)
    dist_rows = [
        ("✔", "Tom Sivil", "Project Manager", "Sales/ Consultancy"),
        ("✔", "Dan White", "Project Manager", "Sales/ Consultancy"),
        ("✔", "Neil Stott", "Director", "Business Development"),
        ("✔", "Ben Harrison", "Managing Director", "N/A"),
        ("✔", "Ellie Sharpe", "Business Dev Exec", "Business Development"),
        ("✔", "Simon Cooper", "Roofing Team Manager", "Installations"),
        ("✔", "Mike Bricknell", "Installation Manager", "Installations"),
        ("✔", "Lauren Symons", "Operations Exec", "Operations"),
        ("✔", "Jack Saunders", "Operations Exec", "Operations"),
        ("✔", "Dan Roote", "Post Installation Manager", "Post Install"),
        ("✔", "Nick Leach", "Electrical Manager", "Electrical"),
        ("✔", "Ashley Pittman", "Electrician", "Electrical"),
        ("✔", "Dan Rogers", "Electrician’s Mate", "Electrical"),
        ("✔", "Nathan Louden", "Installation Team Leader", "Roofing"),
        ("✔", "Tom Attwood", "Installation Team Leader", "Roofing"),
        ("✔", "Installers", "Installation Team Members", "Roofing"),
    ]

    # Light row fills and borders
    for i in range(1, d_rows):
        for j in range(d_cols):
            c = dtable.cell(i,j)
            # Main table background
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*hex_to_rgb("FBE6E1"))
            # Fill hardcoded data when available
            row_idx = i - 1
            if row_idx < len(dist_rows):
                val = dist_rows[row_idx][j]
                c.text_frame.clear()
                para = c.text_frame.paragraphs[0]
                if j == 0:
                    para.alignment = PP_ALIGN.CENTER
                run = para.add_run(); run.text = val
                run.font.color.rgb = RGBColor(0,0,0)
                run.font.size = Pt(9)
            else:
                # Ensure body text color is black for any extra blank rows
                for p in c.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = RGBColor(0,0,0)
                        r.font.size = Pt(9)
            # Borders (matrix grid)
            tc = c._tc.get_or_add_tcPr()
            for tag in ('a:lnL','a:lnR','a:lnT','a:lnB'):
                ln = OxmlElement(tag); ln.set('w','9525')
                sf = OxmlElement('a:solidFill'); clr = OxmlElement('a:srgbClr'); clr.set('val','000000'); sf.append(clr); ln.append(sf)
                tc.append(ln)

def _add_amendments_slide(prs, logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header + logo (logo top-right, title on the left)
    title_text = "Document Name Here"
    title_top = Inches(0.3); title_left = Inches(0.6); title_height = Inches(0.6)
    title_width = prs.slide_width - Inches(1.2)
    if logo_path:
        logo_height = Inches(1.0)
        try:
            slide.shapes.add_picture(str(logo_path), prs.slide_width - Inches(0.6) - logo_height*2.0, Inches(0.2), height=logo_height)
            title_width = prs.slide_width - Inches(1.2) - (logo_height*2.0 + Inches(0.2))
        except Exception:
            pass
    tb = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title_text; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RGBColor(*PALETTE["document_heading"])

    # Section heading
    heading_tb = slide.shapes.add_textbox(Inches(1.9), Inches(1.7), prs.slide_width - Inches(3.8), Inches(0.5))
    htf = heading_tb.text_frame; htf.clear()
    hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run(); hr.text = "Amendments Since Initial Distribution:"; hr.font.size = Pt(16); hr.font.bold = True; hr.font.color.rgb = RGBColor(*PALETTE["document_heading"])

    # Table: 3 columns, 1 header + 10 empty rows
    content_left = Inches(0.6); content_right = prs.slide_width - Inches(0.6)
    available_w = content_right - content_left
    desired_w = Inches(9.0)
    tbl_w = desired_w if desired_w < available_w else available_w
    tbl_left = content_left + (available_w - tbl_w)//2
    tbl_top = Inches(2.2)
    rows, cols = 11, 3
    tbl_h = prs.slide_height - tbl_top - Inches(0.9)
    
    table = slide.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h).table
    # Amendments slide only: make Amendments column the widest
    rev_w = Inches(1.2)
    date_w = Inches(1.5)
    amend_w = tbl_w - (rev_w + date_w)
    if amend_w < Inches(3.0):
        amend_w = Inches(3.0)
        date_w = max(Inches(1.0), tbl_w - (rev_w + amend_w))
    table.columns[0].width = rev_w
    table.columns[1].width = date_w
    table.columns[2].width = amend_w
    headers = ["Revision No", "Date", "Amendments"]
    for j in range(cols):
        hc = table.cell(0,j)
        hc.fill.solid(); hc.fill.fore_color.rgb = RGBColor(*hex_to_rgb("ED1165"))
        hc.text_frame.clear(); hp = hc.text_frame.paragraphs[0]
        hr = hp.add_run(); hr.text = headers[j]; hr.font.bold = True; hr.font.size = Pt(9); hr.font.color.rgb = RGBColor(0,0,0)
        hp.alignment = PP_ALIGN.LEFT if j==2 else PP_ALIGN.CENTER
        tc = hc._tc.get_or_add_tcPr()
        for tag in ('a:lnL','a:lnR','a:lnT','a:lnB'):
            ln = OxmlElement(tag); ln.set('w','9525')
            sf = OxmlElement('a:solidFill'); clr = OxmlElement('a:srgbClr'); clr.set('val','000000'); sf.append(clr); ln.append(sf)
            tc.append(ln)

    for i in range(1, rows):
        for j in range(cols):
            c = table.cell(i,j)
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*hex_to_rgb("FBE6E1"))
            # matrix borders
            tc = c._tc.get_or_add_tcPr()
            for tag in ('a:lnL','a:lnR','a:lnT','a:lnB'):
                ln = OxmlElement(tag); ln.set('w','9525')
                sf = OxmlElement('a:solidFill'); clr = OxmlElement('a:srgbClr'); clr.set('val','000000'); sf.append(clr); ln.append(sf)
                tc.append(ln)
            # default 9pt body
            if c.text_frame.paragraphs:
                for r0 in c.text_frame.paragraphs[0].runs:
                    r0.font.size = Pt(9); r0.font.color.rgb = RGBColor(0,0,0)

    # Footer
    _add_standard_footer(prs, slide)


def _add_notes_slide(prs, logo_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header + logo (logo top-right, title on the left)
    title_text = "Document Name Here"
    title_top = Inches(0.3); title_left = Inches(0.6); title_height = Inches(0.6)
    title_width = prs.slide_width - Inches(1.2)
    if logo_path:
        logo_height = Inches(1.0)
        try:
            slide.shapes.add_picture(str(logo_path), prs.slide_width - Inches(0.6) - logo_height*2.0, Inches(0.2), height=logo_height)
            title_width = prs.slide_width - Inches(1.2) - (logo_height*2.0 + Inches(0.2))
        except Exception:
            pass
    tb = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title_text; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = RGBColor(*PALETTE["document_heading"])

    # Center heading: "Further notes & guidance here"
    heading_tb = slide.shapes.add_textbox(Inches(2.9), Inches(1.2), prs.slide_width - Inches(5.8), Inches(0.5))
    htf = heading_tb.text_frame; htf.clear()
    hp = htf.paragraphs[0]; hp.alignment = PP_ALIGN.CENTER
    hr = hp.add_run(); hr.text = "Further notes & guidance here"; hr.font.size = Pt(20); hr.font.bold = False; hr.font.color.rgb = RGBColor(0,0,0)

    # Large bordered rectangle area (no fill)
    rect_left = Inches(0.3)
    rect_top = Inches(1.6)
    rect_w = prs.slide_width - Inches(0.6)
    rect_h = prs.slide_height - rect_top - Inches(0.9)
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, rect_left, rect_top, rect_w, rect_h)
    try:
        panel.fill.background()
    except Exception:
        panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(255,255,255)
    panel.line.color.rgb = RGBColor(0,0,0)
    panel.line.width = Pt(1.5)
    try:
        panel.shadow.visible = False
        panel.shadow.inherit = False
    except Exception:
        pass

    # Footer: left page count and right document info
    _add_standard_footer(prs, slide)

# ----- Run with external schema file -----
from pathlib import Path

def _detect_base_dir() -> Path:
    """Return the directory to look for input/output files.
    - If running as a frozen executable (PyInstaller), use the executable's directory.
    - Otherwise use the script's directory.
    """
    if getattr(sys, "frozen", False) and hasattr(sys, "executable"):
        return Path(sys.executable).resolve().parent
    return Path(__file__).resolve().parent

def main():
    parser = argparse.ArgumentParser(description="Generate process flow PPTX from a schema .txt")
    parser.add_argument("--schema", type=str, default=None, help="Path to input .txt schema. Defaults to newest .txt in folder.")
    parser.add_argument("--logo", type=str, default=None, help="Optional path to a logo image placed top-right.")
    parser.add_argument("--out", type=str, default=None, help="Optional output .pptx path. Defaults to <schema-stem>.pptx")
    parser.add_argument("--showkey", action="store_true", help="Render a Key box. If omitted, auto-enable when key.txt exists.")
    args = parser.parse_args()

    base_dir = _detect_base_dir()

    # Resolve schema file
    if args.schema:
        schema_file = Path(args.schema)
        if not schema_file.is_absolute():
            schema_file = (base_dir / schema_file).resolve()
        if not schema_file.exists():
            raise FileNotFoundError(f"Schema file not found: {schema_file}")
        if schema_file.name.lower() == "key.txt":
            raise ValueError("key.txt cannot be used as a schema file")
    else:
        txt_files = sorted([
            p for p in base_dir.glob("*.txt")
            if p.is_file() and p.name.lower() != "key.txt"
        ], key=lambda p: p.stat().st_mtime, reverse=True)
        if not txt_files:
            raise FileNotFoundError(f"No .txt files found in {base_dir}")
        schema_file = txt_files[0]

    with open(schema_file, "r", encoding="utf-8") as f:
        schema = f.read()

    # Determine ShowKey flag from schema content
    showkey_in_schema = False
    try:
        # Match lines like: ShowKey: True / true / YES / 1
        m = re.search(r"^\s*ShowKey\s*:\s*(true|1|yes)\s*$", schema, flags=re.IGNORECASE | re.MULTILINE)
        showkey_in_schema = bool(m)
    except Exception:
        showkey_in_schema = False

    # Resolve logo: prefer --logo, else auto-pick my_power_logo.png if present
    logo_path = None
    if args.logo:
        logo_path = Path(args.logo)
        if not logo_path.is_absolute():
            logo_path = (base_dir / logo_path).resolve()
        if not logo_path.exists():
            print(f"Warning: logo not found at {logo_path}; continuing without it.")
            logo_path = None
    else:
        auto_logo = (base_dir / "my_power_logo.png").resolve()
        if auto_logo.exists():
            logo_path = auto_logo

    # Resolve output
    if args.out:
        out_path = Path(args.out)
        if not out_path.is_absolute():
            out_path = (base_dir / out_path).resolve()
    else:
        out_path = (base_dir / f"{schema_file.stem}.pptx").resolve()

    key_path = (base_dir / "key.txt").resolve()
    # Only show key if schema explicitly enables it and key.txt exists
    show_key_effective = showkey_in_schema and key_path.exists()
    render(
        schema,
        str(out_path),
        logo_path=str(logo_path) if logo_path else None,
        show_key=show_key_effective,
        key_path=str(key_path),
    )
    print(f"Read schema from {schema_file.name} -> Saved {out_path}")

if __name__ == "__main__":
    main()
